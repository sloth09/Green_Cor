#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate presentation PPTX for the paper."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.space_after = Pt(6)
    return slide

def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    # Table
    cols = len(headers)
    tbl_rows = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(0.4) * tbl_rows
    table = slide.shapes.add_table(tbl_rows, cols, left, top, width, height).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
    return slide

def add_figure_slide(prs, title, fig_path, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    if os.path.exists(fig_path):
        slide.shapes.add_picture(fig_path, Inches(1.0), Inches(0.9), Inches(8.0), Inches(5.5))
    else:
        txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        txBox2.text_frame.paragraphs[0].text = f"[Figure: {os.path.basename(fig_path)}]"
    if caption:
        txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.5))
        txBox3.text_frame.paragraphs[0].text = caption
        txBox3.text_frame.paragraphs[0].font.size = Pt(10)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    fig_dir = r"D:\code\Green_Cor\results\paper_figures"

    # Slide 1: Title
    add_title_slide(prs,
        "Optimal Ammonia Bunkering Infrastructure\nfor Green Shipping Corridors",
        "A Multi-Period MILP Approach\n\nBusan Port Case Study (2030-2050)")

    # Slide 2: Motivation
    add_content_slide(prs, "Motivation", [
        "IMO 2050: >50% GHG reduction from shipping",
        "144 ammonia-fueled vessels ordered (Dec 2025)",
        "Korea-US green corridor: Busan-Tacoma by 2027",
        "$10B Busan mega-port investment",
        "Missing: port-level bunkering infrastructure sizing"
    ])

    # Slide 3: Research Gaps
    add_content_slide(prs, "Research Gaps", [
        "Gap 1: No joint optimization of shuttle-pump-fleet for ammonia",
        "Gap 2: No quantitative port storage vs remote supply comparison",
        "Gap 3: No multi-period fleet expansion model for ammonia bunkering",
        "Closest work: Yang & Lam (2023) DES -- simulation, not optimization"
    ])

    # Slide 4: Methodology
    add_content_slide(prs, "Methodology: MILP Model", [
        "Objective: Minimize 20-year Net Present Cost (NPC)",
        "Decision variables: shuttle size, pump rate, fleet additions per year",
        "Constraints: demand, working time (8,000 h/yr), tank capacity",
        "3 supply chain configurations: Busan / Yeosu(86nm) / Ulsan(59nm)",
        "Cost: CAPEX scaling law (alpha=0.75), OPEX (fuel, maintenance)",
        "Annuity factor: AF = 10.8355 (r=7%, n=21 years)"
    ])

    # Slide 5: Three Cases
    add_table_slide(prs, "Three Supply Chain Configurations",
        ["", "Case 1 (Busan)", "Case 2 (Ulsan)", "Case 3 (Yeosu)"],
        [
            ["Source", "Port storage", "Ulsan (59 nm)", "Yeosu (86 nm)"],
            ["Travel time", "1.0 h", "3.93 h", "5.73 h"],
            ["Shuttle range", "500-10,000 m3", "2,500-50,000 m3", "2,500-50,000 m3"],
            ["Pumping logic", "Shuttle/pump", "Vessel/pump", "Vessel/pump"],
            ["Storage", "35,000 ton tank", "None at Busan", "None at Busan"],
        ])

    # Slide 6: NPC vs Shuttle Size (D1)
    add_figure_slide(prs, "Fig. 2: NPC vs Shuttle Size",
        os.path.join(fig_dir, "D1_npc_vs_shuttle.png"),
        "Convex NPC curves with distinct optima: 2,500 / 10,000 / 5,000 m3")

    # Slide 7: Optimal Results Table
    add_table_slide(prs, "Optimal Configurations (Key Results)",
        ["", "Case 1", "Case 2", "Case 3"],
        [
            ["Shuttle (m3)", "2,500", "5,000", "10,000"],
            ["NPC (USD M)", "290.81", "700.68", "879.88"],
            ["LCOA ($/ton)", "1.23", "2.97", "3.73"],
            ["Cycle time (h)", "10.17", "23.19", "38.13"],
            ["Ratio vs Case 1", "1.0x", "2.41x", "3.02x"],
        ])

    # Slide 8: Cost Breakdown (D6)
    add_figure_slide(prs, "Fig. 4: Cost Component Breakdown",
        os.path.join(fig_dir, "D6_cost_breakdown.png"),
        "Case 1: CAPEX-dominant (45.6%); Case 2: vOPEX-dominant (33-39%)")

    # Slide 9: Fleet Evolution (D8)
    add_figure_slide(prs, "Fig. 7: Fleet Size Evolution",
        os.path.join(fig_dir, "D8_fleet_evolution.png"),
        "Staircase fleet expansion synchronized with linear demand growth")

    # Slide 10: Tornado (FIG7)
    add_figure_slide(prs, "Fig. 10: Tornado Sensitivity Analysis",
        os.path.join(fig_dir, "Fig7_tornado_deterministic.png"),
        "Case 1: CAPEX Scaling dominates (62%); Case 2: Bunker Volume dominates")

    # Slide 11: Fuel Price (FIG8)
    add_figure_slide(prs, "Fig. 11: Fuel Price Sensitivity",
        os.path.join(fig_dir, "Fig8_fuel_price_sensitivity.png"),
        "Case 1 NPC: $255M-$362M across $300-$1,200/ton range")

    # Slide 12: Break-even (FIG9)
    add_figure_slide(prs, "Fig. 12: Break-Even Distance Analysis",
        os.path.join(fig_dir, "Fig9_breakeven_distance.png"),
        "Yeosu: crossover ~59.6 nm; Ulsan: no crossover (Case 1 always cheaper)")

    # Slide 13: Demand Scenarios (FIG10)
    add_figure_slide(prs, "Fig. 13: Demand Scenario Analysis",
        os.path.join(fig_dir, "Fig10_demand_scenarios.png"),
        "LCOA stable: $1.21-$1.28/ton across 4x demand range (5.7% variation)")

    # Slide 14: Key Findings
    add_content_slide(prs, "Key Findings", [
        "C1: Optimal shuttle = 2,500 m3 (Case 1), NPC $290.81M",
        "C2: Break-even distance ~59.6 nm for port storage vs remote supply",
        "C3: Shuttle specs invariant to 4x demand range (LCO range 5.7%)",
        "C4: Cost drivers differ: CAPEX scaling (Case 1) vs Bunker volume (Case 2)"
    ])

    # Slide 15: Practical Recommendations
    add_content_slide(prs, "Practical Recommendations", [
        "1. Build port-based storage for distances > 60 nm",
        "2. Commit to shuttle specifications early (robust to demand uncertainty)",
        "3. Monitor shipyard costs as primary risk factor (62% NPC swing)",
        "4. LCOA of $1.23/ton = 0.1-0.4% of total ammonia fuel cost"
    ])

    # Slide 16: Limitations & Future Work
    add_content_slide(prs, "Limitations & Future Work", [
        "Limitations: linear demand, fixed fuel price, no discounting, no queuing",
        "Future: stochastic MILP, DES-MILP hybrid, multi-fuel comparison",
        "Future: real-options analysis, multi-port network extension",
        "Framework transferable to other ports (substitute local parameters)"
    ])

    # Slide 17: Thank You
    add_title_slide(prs, "Thank You", "Questions & Discussion")

    out_path = r"D:\code\Green_Cor\docs\paper\presentation.pptx"
    prs.save(out_path)
    print(f"[OK] Presentation saved to {out_path}")
    print(f"[OK] Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
